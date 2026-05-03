# core/bp_parser.py
import json
import re
import anthropic

def extract_text_from_pdf(file_path: str) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception:
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts)
    except Exception:
        return ""

def extract_text_from_file(file_path: str) -> str:
    path = file_path.lower()
    if path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif path.endswith((".pptx", ".ppt")):
        return extract_text_from_pptx(file_path)
    return ""

EXTRACT_PROMPT = """你是一名专业的投资 FA 助手。请从以下 BP（商业计划书）文本中提取关键信息，以 JSON 格式返回。

BP 文本：
{text}

请提取以下字段，如果信息不存在则返回空字符串：
- name: 项目/公司名称
- sector: 所在赛道（如 AI、消费、医疗健康、企业服务、硬件等）
- sub_sector: 细分领域（更具体的方向，如 AI+工业质检、新能源储能等）
- stage: 当前融资阶段（天使轮、Pre-A轮、A轮、B轮等）
- location: 公司所在地（城市）
- description: 项目简介（100字以内）
- highlights: 核心亮点（用分号分隔，列出3-5点）
- financing_need: 本轮融资需求金额

只返回 JSON，不要其他内容。"""

def extract_project_info(text: str, api_key: str) -> dict:
    client = anthropic.Anthropic(api_key=api_key)
    response = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1024,
        messages=[{"role": "user", "content": EXTRACT_PROMPT.format(text=text[:8000])}]
    )
    raw = response.content[0].text.strip()
    raw = re.sub(r"^```json\s*|\s*```$", "", raw, flags=re.MULTILINE).strip()
    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        data = {}
    defaults = {"name": "", "sector": "", "sub_sector": "", "stage": "",
                "location": "", "description": "", "highlights": "", "financing_need": ""}
    return {**defaults, **data}
