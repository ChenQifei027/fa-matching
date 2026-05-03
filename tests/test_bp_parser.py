# tests/test_bp_parser.py
import pytest
from unittest.mock import patch, MagicMock
from core.bp_parser import extract_text_from_pdf, extract_text_from_pptx, extract_project_info

def test_extract_text_from_pdf(tmp_path):
    fake_pdf = tmp_path / "test.pdf"
    fake_pdf.write_bytes(b"")
    result = extract_text_from_pdf(str(fake_pdf))
    assert isinstance(result, str)

def test_extract_text_from_pptx(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "测试项目"
    slide.placeholders[1].text = "我们是一家AI公司，专注于工业质检。"
    pptx_path = tmp_path / "test.pptx"
    prs.save(str(pptx_path))
    text = extract_text_from_pptx(str(pptx_path))
    assert "测试项目" in text
    assert "工业质检" in text

def test_extract_project_info_returns_required_fields(mocker):
    mock_client = mocker.patch("core.bp_parser.anthropic.Anthropic")
    mock_instance = MagicMock()
    mock_client.return_value = mock_instance
    mock_instance.messages.create.return_value = MagicMock(
        content=[MagicMock(text='{"name": "智检科技", "sector": "AI", "sub_sector": "AI+工业质检", "stage": "Pre-A", "location": "深圳", "description": "基于AI的工业质检系统", "highlights": "准确率99.5%，已服务50家工厂", "financing_need": "3000万人民币"}')]
    )
    result = extract_project_info("这是一段BP文本内容", api_key="test_key")
    assert result["sector"] == "AI"
    assert result["sub_sector"] == "AI+工业质检"
    assert result["stage"] == "Pre-A"
    assert "name" in result
    assert "description" in result

def test_extract_project_info_handles_partial_extraction(mocker):
    mock_client = mocker.patch("core.bp_parser.anthropic.Anthropic")
    mock_instance = MagicMock()
    mock_client.return_value = mock_instance
    mock_instance.messages.create.return_value = MagicMock(
        content=[MagicMock(text='{"name": "某项目", "sector": "AI"}')]
    )
    result = extract_project_info("短文本", api_key="test_key")
    assert result.get("sub_sector", "") == ""
    assert result["sector"] == "AI"
